"""Bounded, path-free previews for uploaded PDF, PPTX, and XLSX documents."""

from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Any, Mapping
from zipfile import BadZipFile, ZipFile

import fitz
import openpyxl
from pptx import Presentation

from storage import MAX_FILE_BYTES


MAX_PREVIEW_UNITS = 3
MAX_PREVIEW_ROWS = 20
MAX_PREVIEW_COLUMNS = 12
MAX_PREVIEW_TEXT = 8_000
MAX_ARCHIVE_ENTRIES = 2_000
MAX_ARCHIVE_EXPANDED_BYTES = 100 * 1024 * 1024


class DocumentPreviewError(ValueError):
    """Raised with a display-safe message when a preview cannot be produced."""


def _bounded_office_archive(contents: bytes) -> None:
    try:
        with ZipFile(BytesIO(contents)) as archive:
            entries = archive.infolist()
            if (
                len(entries) > MAX_ARCHIVE_ENTRIES
                or sum(entry.file_size for entry in entries) > MAX_ARCHIVE_EXPANDED_BYTES
            ):
                raise DocumentPreviewError("문서가 미리보기 안전 한도를 초과했습니다.")
    except BadZipFile as error:
        raise DocumentPreviewError("문서 미리보기를 만들 수 없습니다.") from error


def _text(value: object) -> str:
    return str(value)[:MAX_PREVIEW_TEXT]


def build_document_preview(name: str, contents: bytes) -> dict[str, Any]:
    """Return a bounded preview without retaining or exposing a server path."""

    if not isinstance(name, str) or not isinstance(contents, bytes):
        raise DocumentPreviewError("문서 미리보기를 만들 수 없습니다.")
    if not contents or len(contents) > MAX_FILE_BYTES:
        raise DocumentPreviewError("문서가 미리보기 안전 한도를 초과했습니다.")
    suffix = Path(name).suffix.lower()
    try:
        if suffix == ".pdf":
            document = fitz.open(stream=contents, filetype="pdf")
            try:
                page_count = len(document)
                pages = []
                for number in range(min(page_count, MAX_PREVIEW_UNITS)):
                    page = document[number]
                    largest = max(float(page.rect.width), float(page.rect.height), 1.0)
                    scale = min(1.5, 1600.0 / largest)
                    pixmap = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
                    pages.append({"label": f"페이지 {number + 1}", "image": pixmap.tobytes("png")})
            finally:
                document.close()
            return {"kind": "pdf", "units": pages, "truncated": page_count > MAX_PREVIEW_UNITS}

        if suffix == ".pptx":
            _bounded_office_archive(contents)
            presentation = Presentation(BytesIO(contents))
            slides = []
            for number, slide in enumerate(presentation.slides, start=1):
                if number > MAX_PREVIEW_UNITS:
                    break
                texts = []
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False) and shape.text.strip():
                        texts.append(shape.text.strip())
                    if getattr(shape, "has_table", False):
                        for row_number, row in enumerate(shape.table.rows):
                            if row_number >= MAX_PREVIEW_ROWS:
                                break
                            texts.append(" | ".join(
                                cell.text.strip() for column, cell in enumerate(row.cells)
                                if column < MAX_PREVIEW_COLUMNS
                            ))
                slides.append({"label": f"슬라이드 {number}", "text": _text("\n".join(texts))})
            return {"kind": "pptx", "units": slides, "truncated": len(presentation.slides) > MAX_PREVIEW_UNITS}

        if suffix == ".xlsx":
            _bounded_office_archive(contents)
            workbook = openpyxl.load_workbook(
                BytesIO(contents), read_only=True, data_only=True, keep_links=False,
            )
            try:
                sheets = []
                for worksheet in workbook.worksheets[:MAX_PREVIEW_UNITS]:
                    rows = []
                    for row_number, row in enumerate(worksheet.iter_rows(values_only=True)):
                        if row_number >= MAX_PREVIEW_ROWS:
                            break
                        rows.append([_text("" if value is None else value) for value in row[:MAX_PREVIEW_COLUMNS]])
                    sheets.append({"label": _text(worksheet.title), "rows": rows})
            finally:
                workbook.close()
            return {"kind": "xlsx", "units": sheets, "truncated": len(workbook.sheetnames) > MAX_PREVIEW_UNITS}
    except DocumentPreviewError:
        raise
    except Exception as error:
        raise DocumentPreviewError("문서 미리보기를 만들 수 없습니다.") from error
    raise DocumentPreviewError("지원하지 않는 문서 형식입니다.")


def render_document_preview(st: object, preview: Mapping[str, object]) -> None:
    """Render only the bounded preview structure returned above."""

    units = preview.get("units", [])
    if not isinstance(units, list) or not units:
        st.info("표시할 미리보기 내용이 없습니다.")
        return
    for unit in units:
        if not isinstance(unit, Mapping):
            continue
        st.caption(str(unit.get("label", "미리보기")))
        if preview.get("kind") == "pdf" and isinstance(unit.get("image"), bytes):
            st.image(unit["image"], use_container_width=True)
        elif preview.get("kind") == "xlsx":
            st.dataframe(unit.get("rows", []), use_container_width=True, hide_index=True)
        else:
            st.text(str(unit.get("text", "")))
    if preview.get("truncated") is True:
        st.caption("안전한 미리보기를 위해 일부 페이지만 표시합니다.")
