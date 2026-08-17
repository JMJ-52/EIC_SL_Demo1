from __future__ import annotations

import os
from collections.abc import Iterable, Iterator

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from lifecycle.extraction.models import ParsedUnit


def _iter_shapes(shapes: Iterable) -> Iterator:
    """Yield shapes, recursing into groups so grouped content isn't lost."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def parse_pptx(path: str) -> list[ParsedUnit]:
    source_doc = os.path.basename(path)
    presentation = Presentation(path)
    units = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        text_parts: list[str] = []
        images: list[bytes] = []
        for shape in _iter_shapes(slide.shapes):
            if shape.has_text_frame and shape.text_frame.text.strip():
                text_parts.append(shape.text_frame.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        text_parts.append(row_text)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                images.append(shape.image.blob)
        units.append(
            ParsedUnit(source_doc=source_doc, unit_label=f"슬라이드 {slide_number}",
                       text="\n".join(text_parts), images=images)
        )
    return units

