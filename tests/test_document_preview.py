from io import BytesIO

import fitz
from openpyxl import Workbook
from pptx import Presentation

from document_preview import build_document_preview


def _pdf_bytes(page_count: int = 1) -> bytes:
    document = fitz.open()
    for number in range(page_count):
        page = document.new_page()
        page.insert_text((72, 72), f"page {number + 1}")
    contents = document.tobytes()
    document.close()
    return contents


def _pptx_bytes() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Replacement review"
    slide.placeholders[1].text = "Original slide content"
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _xlsx_bytes() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Inspection"
    sheet.append(["Equipment", "Status"])
    sheet.append(["Motor A", "Reviewed"])
    output = BytesIO()
    workbook.save(output)
    workbook.close()
    return output.getvalue()


def test_pdf_preview_renders_only_bounded_original_pages() -> None:
    preview = build_document_preview("source.pdf", _pdf_bytes(4))

    assert preview["kind"] == "pdf"
    assert len(preview["units"]) == 3
    assert preview["units"][0]["image"].startswith(b"\x89PNG")
    assert preview["truncated"] is True


def test_pptx_preview_extracts_bounded_original_slide_content() -> None:
    preview = build_document_preview("source.pptx", _pptx_bytes())

    assert preview["kind"] == "pptx"
    assert "Replacement review" in preview["units"][0]["text"]
    assert "Original slide content" in preview["units"][0]["text"]


def test_xlsx_preview_extracts_bounded_original_sheet_rows() -> None:
    preview = build_document_preview("source.xlsx", _xlsx_bytes())

    assert preview["kind"] == "xlsx"
    assert preview["units"][0]["label"] == "Inspection"
    assert preview["units"][0]["rows"][1] == ["Motor A", "Reviewed"]
